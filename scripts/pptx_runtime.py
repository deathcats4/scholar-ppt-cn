from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from scripts.common import Issue

try:
    from PIL import Image as PillowImage
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PillowImage = Presentation = RGBColor = None
    MSO_CONNECTOR = MSO_SHAPE = MSO_ANCHOR = PP_ALIGN = None
    OxmlElement = qn = Inches = Pt = None
    PPTX_AVAILABLE = False


DEFAULT_COLORS = {
    "background": "FFFFFF",
    "primary": "123B86",
    "secondary": "4C6FAE",
    "text": "1E2530",
    "muted": "687386",
    "line": "D9E3F3",
    "panel": "F5F8FC",
}
DEFAULT_CHINESE_FONT = "Microsoft YaHei"
DEFAULT_LATIN_FONT = "Times New Roman"
_CJK_RE = re.compile(r"[\u3000-\u303f\u3400-\u9fff\uf900-\ufaff\uff00-\uffef]")
_INTERNAL_RENDER_TYPES = {
    "cover",
    "section",
    "statement",
    "bullets",
    "figure",
    "comparison",
    "multi-panel",
    "process",
    "conclusion",
}
RENDER_CONTENT_LIMITS = {
    "cover": {"items": 4},
    "statement": {"items": 4},
    "bullets": {"items": 4},
    "figure": {"items": 4, "body": 4},
    "process": {"items": 5, "body": 5},
    "conclusion": {"items": 4, "body": 4},
}


@dataclass(frozen=True)
class Theme:
    background: str
    primary: str
    secondary: str
    text: str
    muted: str
    line: str
    panel: str
    chinese_title_font: str
    chinese_body_font: str
    latin_font: str


@dataclass
class BuildContext:
    presentation: Any
    project: dict[str, Any]
    base_dir: Path
    width: float
    height: float
    theme: Theme
    assets: dict[str, dict[str, Any]]
    issues: list[Issue]
    rendered_asset_ids: set[str]


def _hex(value: str | None, fallback: str) -> str:
    candidate = str(value or "").lstrip("#")
    return candidate.upper() if re.fullmatch(r"[0-9A-Fa-f]{6}", candidate) else fallback


def _first_font(project: dict[str, Any], role: str, fallback: str) -> str:
    try:
        families = project["template"]["dna"]["font_roles"][role]["families"]
    except (KeyError, TypeError):
        return fallback
    if not isinstance(families, list):
        return fallback
    selected = next((item for item in families if isinstance(item, str) and item), None)
    return selected or fallback


def theme_from_project(project: dict[str, Any]) -> Theme:
    colors = project.get("template", {}).get("dna", {}).get("colors", [])
    colors = colors if isinstance(colors, list) else []
    primary = _hex(colors[0] if colors else None, DEFAULT_COLORS["primary"])
    secondary = _hex(colors[1] if len(colors) > 1 else None, DEFAULT_COLORS["secondary"])
    return Theme(
        background=DEFAULT_COLORS["background"],
        primary=primary,
        secondary=secondary,
        text=DEFAULT_COLORS["text"],
        muted=DEFAULT_COLORS["muted"],
        line=DEFAULT_COLORS["line"],
        panel=DEFAULT_COLORS["panel"],
        chinese_title_font=_first_font(project, "title", DEFAULT_CHINESE_FONT),
        chinese_body_font=_first_font(project, "body", DEFAULT_CHINESE_FONT),
        latin_font=DEFAULT_LATIN_FONT,
    )


def _rgb(value: str) -> Any:
    return RGBColor.from_string(value)


def _segments(text: str) -> Iterable[tuple[str, bool]]:
    if not text:
        return
    start = 0
    current = bool(_CJK_RE.match(text[0]))
    for index, character in enumerate(text[1:], start=1):
        is_cjk = bool(_CJK_RE.match(character))
        if is_cjk != current:
            yield text[start:index], current
            start = index
            current = is_cjk
    yield text[start:], current


def _set_east_asian_font(run: Any, font_name: str) -> None:
    run.font.name = font_name
    properties = run._r.get_or_add_rPr()
    element = properties.find(qn("a:ea"))
    if element is None:
        element = OxmlElement("a:ea")
        properties.append(element)
    element.set("typeface", font_name)


def _add_runs(
    paragraph: Any,
    text: str,
    theme: Theme,
    *,
    size: float,
    color: str,
    bold: bool = False,
    font_role: str = "body",
) -> None:
    chinese_font = (
        theme.chinese_title_font
        if font_role == "title"
        else theme.chinese_body_font
    )
    for segment, cjk in _segments(text):
        run = paragraph.add_run()
        run.text = segment
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        _set_east_asian_font(
            run,
            chinese_font if cjk else theme.latin_font,
        )


def _add_text(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    theme: Theme,
    *,
    size: float,
    color: str | None = None,
    bold: bool = False,
    align: str = "left",
    valign: str = "top",
    margin: float = 0.02,
    font_role: str = "body",
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    _add_runs(
        paragraph,
        str(text),
        theme,
        size=size,
        color=color or theme.text,
        bold=bold,
        font_role=font_role,
    )
    return box


def _add_paragraphs(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    paragraphs: list[str],
    theme: Theme,
    *,
    size: float = 18,
    bullet: bool = False,
    color: str | None = None,
    spacing: float = 9,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, text in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(spacing)
        display = f"•  {text}" if bullet else text
        _add_runs(
            paragraph,
            display,
            theme,
            size=size,
            color=color or theme.text,
        )
    return box


def _add_rect(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1,
    rounded: bool = False,
) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def _add_line(
    slide: Any,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str,
    width: float = 1.2,
) -> Any:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _add_page_header(
    slide: Any,
    context: BuildContext,
    slide_data: dict[str, Any],
    *,
    centered: bool = False,
) -> None:
    section = str(slide_data.get("narrative_section") or "").strip()
    if section:
        _add_rect(
            slide,
            0.42,
            0.30,
            0.08,
            0.32,
            fill=context.theme.primary,
        )
        _add_text(
            slide,
            0.60,
            0.26,
            2.8,
            0.42,
            section,
            context.theme,
            size=15,
            color=context.theme.primary,
            bold=True,
            valign="middle",
        )
    title_x = 0.68 if not centered else 1.0
    title_w = context.width - (1.36 if not centered else 2.0)
    _add_text(
        slide,
        title_x,
        0.78,
        title_w,
        0.72,
        str(slide_data.get("title") or ""),
        context.theme,
        size=28,
        color=context.theme.primary,
        bold=True,
        align="center" if centered else "left",
        valign="middle",
        font_role="title",
    )
    if centered:
        _add_line(
            slide,
            context.width / 2 - 0.35,
            1.54,
            context.width / 2 + 0.35,
            1.54,
            context.theme.primary,
            1.6,
        )


def _add_page_number(slide: Any, context: BuildContext, number: int) -> None:
    if number <= 1:
        return
    _add_text(
        slide,
        context.width - 0.62,
        context.height - 0.42,
        0.28,
        0.18,
        str(number),
        context.theme,
        size=9,
        color=context.theme.muted,
        align="right",
    )


def _resolve_asset_path(context: BuildContext, asset: dict[str, Any]) -> Path | None:
    value = asset.get("path")
    if not isinstance(value, str) or not value.strip():
        return None
    path = Path(value)
    if path.is_absolute():
        return path
    return (context.base_dir / path).resolve()


def _asset_ids(slide_data: dict[str, Any], explicit_only: bool = False) -> list[str]:
    render = slide_data.get("render")
    if (
        isinstance(render, dict)
        and "asset_ids" in render
        and isinstance(render.get("asset_ids"), list)
    ):
        values = [item for item in render["asset_ids"] if isinstance(item, str)]
        return values
    if explicit_only:
        return []
    values = slide_data.get("source_asset_ids", [])
    return [item for item in values if isinstance(item, str)] if isinstance(values, list) else []


def _add_image_contain(
    slide: Any,
    context: BuildContext,
    asset_id: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str = "",
) -> bool:
    asset = context.assets.get(asset_id)
    if asset is None:
        context.issues.append(
            Issue("error", "build.asset_reference", f"Unknown asset: {asset_id}", asset_id)
        )
        return False
    path = _resolve_asset_path(context, asset)
    if path is None or not path.is_file():
        context.issues.append(
            Issue(
                "error",
                "build.asset_missing",
                f"Asset file does not exist: {asset.get('path')!r}",
                asset_id,
            )
        )
        return False
    try:
        with PillowImage.open(path) as image:
            pixel_width, pixel_height = image.size
    except (OSError, ValueError) as exc:
        context.issues.append(
            Issue("error", "build.asset_decode", f"Cannot decode image: {exc}", asset_id)
        )
        return False
    if pixel_width <= 0 or pixel_height <= 0:
        context.issues.append(
            Issue("error", "build.asset_decode", "Image has invalid dimensions", asset_id)
        )
        return False
    ratio = pixel_width / pixel_height
    box_ratio = w / h
    if ratio >= box_ratio:
        draw_w = w
        draw_h = w / ratio
    else:
        draw_h = h
        draw_w = h * ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    effective_ppi = min(pixel_width / draw_w, pixel_height / draw_h)
    if effective_ppi < 120:
        context.issues.append(
            Issue(
                "warning",
                "build.asset_resolution",
                f"Effective image resolution is approximately {effective_ppi:.0f} ppi",
                asset_id,
            )
        )
    _add_rect(slide, x, y, w, h, fill="FFFFFF", line=context.theme.line)
    try:
        slide.shapes.add_picture(
            str(path),
            Inches(draw_x),
            Inches(draw_y),
            width=Inches(draw_w),
            height=Inches(draw_h),
        )
    except Exception as exc:
        context.issues.append(
            Issue(
                "error",
                "build.asset_insert",
                f"Cannot insert image into PPTX: {type(exc).__name__}: {exc}",
                asset_id,
            )
        )
        return False
    context.rendered_asset_ids.add(asset_id)
    if label:
        _add_text(
            slide,
            x,
            y + h + 0.05,
            w,
            0.28,
            label,
            context.theme,
            size=11,
            color=context.theme.muted,
            align="center",
        )
    return True


def _render(slide_data: dict[str, Any]) -> dict[str, Any]:
    value = slide_data.get("render")
    return value if isinstance(value, dict) else {}


def _render_items(slide_data: dict[str, Any]) -> list[dict[str, Any]]:
    items = _render(slide_data).get("items", [])
    return [item for item in items if isinstance(item, dict)] if isinstance(items, list) else []


def _items_for_assets(
    slide_data: dict[str, Any],
    asset_ids: list[str],
) -> list[dict[str, Any]]:
    items = _render_items(slide_data)
    by_asset = {
        item["asset_id"]: item
        for item in items
        if isinstance(item.get("asset_id"), str)
    }
    bodies = _body(slide_data)
    resolved: list[dict[str, Any]] = []
    for index, asset_id in enumerate(asset_ids):
        item = dict(by_asset.get(asset_id, {}))
        if not item and index < len(items) and items[index].get("asset_id") is None:
            item = dict(items[index])
        if not str(item.get("body") or "").strip() and index < len(bodies):
            item["body"] = bodies[index]
        item["asset_id"] = asset_id
        resolved.append(item)
    return resolved


def infer_render_type(slide_data: dict[str, Any]) -> str:
    render_type = _render(slide_data).get("type")
    if isinstance(render_type, str) and render_type in _INTERNAL_RENDER_TYPES:
        return render_type
    number = slide_data.get("number")
    title = str(slide_data.get("title") or "")
    section = str(slide_data.get("narrative_section") or "")
    asset_count = len(_asset_ids(slide_data))
    if number == 1 or "封面" in section:
        return "cover"
    if asset_count >= 3:
        return "multi-panel"
    if asset_count == 2:
        return "comparison"
    if asset_count == 1:
        return "figure"
    if any(keyword in title + section for keyword in ("结论", "总结", "展望")):
        return "conclusion"
    task = str(slide_data.get("communication_task") or "")
    if any(keyword in title + task for keyword in ("流程", "步骤", "机制", "路径")):
        return "process"
    return "bullets"


def _body(slide_data: dict[str, Any]) -> list[str]:
    body = _render(slide_data).get("body")
    if isinstance(body, list):
        result = [item.strip() for item in body if isinstance(item, str) and item.strip()]
        if result:
            return result
    core = str(slide_data.get("core_message") or "").strip()
    return [core] if core else []


def _render_cover(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    render = _render(data)
    explicit_assets = (
        _asset_ids(data)
        if render.get("type") == "cover"
        else _asset_ids(data, explicit_only=True)
    )
    has_visual = bool(explicit_assets)
    title_w = context.width * (0.56 if has_visual else 0.80)
    _add_text(
        slide,
        0.72,
        0.82,
        3.2,
        0.42,
        str(data.get("narrative_section") or "学术汇报"),
        context.theme,
        size=16,
        color=context.theme.primary,
        bold=True,
    )
    _add_text(
        slide,
        0.72,
        1.58,
        title_w,
        2.10,
        str(data.get("title") or ""),
        context.theme,
        size=34,
        color=context.theme.primary,
        bold=True,
        valign="middle",
        font_role="title",
    )
    subtitle = render.get("subtitle") or data.get("core_message")
    if isinstance(subtitle, str) and subtitle.strip() and subtitle.strip() != data.get("title"):
        _add_line(slide, 0.72, 4.02, 1.35, 4.02, context.theme.primary, 1.6)
        _add_text(
            slide,
            1.48,
            3.78,
            title_w - 0.76,
            0.52,
            subtitle,
            context.theme,
            size=18,
            color=context.theme.text,
            valign="middle",
        )
    items = _render_items(data)
    metadata = [
        "　".join(
            part
            for part in (
                str(item.get("title") or "").strip(),
                str(item.get("body") or "").strip(),
            )
            if part
        )
        for item in items
    ]
    metadata = [item for item in metadata if item]
    if metadata:
        _add_paragraphs(
            slide,
            0.72,
            4.72,
            min(title_w, 6.5),
            1.15,
            metadata[:4],
            context.theme,
            size=15,
            spacing=5,
        )
    if has_visual:
        _add_image_contain(
            slide,
            context,
            explicit_assets[0],
            context.width * 0.65,
            1.22,
            context.width * 0.29,
            context.height * 0.70,
        )


def _render_section(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_rect(slide, 0.72, 0.72, 0.10, 0.52, fill=context.theme.primary)
    _add_text(
        slide,
        1.02,
        0.68,
        2.8,
        0.56,
        str(data.get("narrative_section") or "章节"),
        context.theme,
        size=18,
        color=context.theme.primary,
        bold=True,
        valign="middle",
    )
    _add_text(
        slide,
        1.05,
        2.05,
        context.width - 2.1,
        1.55,
        str(data.get("title") or ""),
        context.theme,
        size=36,
        color=context.theme.primary,
        bold=True,
        align="center",
        valign="middle",
        font_role="title",
    )
    subtitle = _render(data).get("subtitle") or data.get("core_message")
    if isinstance(subtitle, str) and subtitle.strip() and subtitle.strip() != data.get("title"):
        _add_text(
            slide,
            2.2,
            4.1,
            context.width - 4.4,
            0.68,
            subtitle,
            context.theme,
            size=18,
            color=context.theme.muted,
            align="center",
            valign="middle",
        )


def _render_bullets(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_page_header(slide, context, data)
    items = _render_items(data)
    if items:
        row_height = min(1.2, 4.6 / max(1, len(items)))
        y = 1.72
        for index, item in enumerate(items[:4], start=1):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.82),
                Inches(y - 0.02),
                Inches(0.52),
                Inches(0.52),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = _rgb(context.theme.primary)
            circle.line.fill.background()
            _add_text(
                slide,
                0.86,
                y,
                0.50,
                0.48,
                str(index),
                context.theme,
                size=16,
                color="FFFFFF",
                bold=True,
                align="center",
                valign="middle",
            )
            _add_text(
                slide,
                1.55,
                y - 0.05,
                context.width - 2.35,
                0.42,
                str(item.get("title") or ""),
                context.theme,
                size=20,
                color=context.theme.primary,
                bold=True,
            )
            body = str(item.get("body") or "").strip()
            if body:
                _add_text(
                    slide,
                    1.55,
                    y + 0.40,
                    context.width - 2.35,
                    0.42,
                    body,
                    context.theme,
                    size=14,
                    color=context.theme.text,
                )
            y += row_height
    else:
        body = _body(data)
        if body:
            _add_paragraphs(
                slide,
                1.0,
                2.0,
                context.width - 2.0,
                3.7,
                body,
                context.theme,
                size=22 if len(body) == 1 else 18,
                bullet=len(body) > 1,
                spacing=14,
            )


def _render_figure(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_page_header(slide, context, data)
    ids = _asset_ids(data)
    figure_w = context.width * 0.66
    if ids:
        _add_image_contain(slide, context, ids[0], 0.62, 1.72, figure_w, 4.85)
    else:
        context.issues.append(
            Issue(
                "warning",
                "build.figure_without_asset",
                "Figure slide has no source asset; rendered as a text slide",
                str(data.get("id") or ""),
            )
        )
    x = 0.62 + figure_w + 0.36
    items = _render_items(data)
    body = [
        str(item.get("title") or item.get("body") or "").strip()
        for item in items
        if str(item.get("title") or item.get("body") or "").strip()
    ] or _body(data)
    if body:
        _add_rect(
            slide,
            x,
            1.72,
            context.width - x - 0.62,
            4.85,
            fill=context.theme.panel,
            line=context.theme.line,
            rounded=True,
        )
        _add_paragraphs(
            slide,
            x + 0.24,
            2.02,
            context.width - x - 1.10,
            4.2,
            body[:4],
            context.theme,
            size=16,
            bullet=len(body) > 1,
            spacing=12,
        )


def _render_comparison(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_page_header(slide, context, data, centered=True)
    ids = _asset_ids(data)
    items = _items_for_assets(data, ids)
    gap = 0.34
    panel_w = (context.width - 1.36 - gap) / 2
    for index in range(2):
        x = 0.68 + index * (panel_w + gap)
        title = (
            str(items[index].get("title") or "")
            if index < len(items)
            else f"对比对象 {index + 1}"
        )
        if index < len(ids):
            _add_image_contain(
                slide,
                context,
                ids[index],
                x,
                1.88,
                panel_w,
                3.92,
                label=title,
            )
        else:
            _add_rect(
                slide,
                x,
                1.88,
                panel_w,
                3.92,
                fill=context.theme.panel,
                line=context.theme.line,
            )
            _add_text(
                slide,
                x + 0.3,
                3.1,
                panel_w - 0.6,
                0.75,
                title,
                context.theme,
                size=20,
                color=context.theme.primary,
                bold=True,
                align="center",
                valign="middle",
            )
        if index < len(items) and str(items[index].get("body") or "").strip():
            _add_text(
                slide,
                x,
                6.18,
                panel_w,
                0.45,
                str(items[index]["body"]),
                context.theme,
                size=13,
                color=context.theme.text,
                align="center",
            )


def _render_multi_panel(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_page_header(slide, context, data)
    ids = _asset_ids(data)
    items = _items_for_assets(data, ids)
    columns = 2
    rows = 2 if len(ids) > 2 else 1
    gap_x, gap_y = 0.28, 0.32
    panel_w = (context.width - 1.36 - gap_x) / columns
    panel_h = (4.95 - gap_y * (rows - 1)) / rows
    for index, asset_id in enumerate(ids):
        row, column = divmod(index, columns)
        x = 0.68 + column * (panel_w + gap_x)
        y = 1.72 + row * (panel_h + gap_y)
        label = str(items[index].get("title") or "") if index < len(items) else ""
        _add_image_contain(
            slide,
            context,
            asset_id,
            x,
            y,
            panel_w,
            panel_h - (0.28 if label else 0),
            label=label,
        )
    if not ids:
        body = _body(data)
        if body:
            _add_text(
                slide,
                0.82,
                1.80,
                context.width - 1.64,
                4.85,
                "\n".join(f"• {text}" for text in body),
                context.theme,
                size=22,
                color=context.theme.text,
                valign="top",
            )


def _render_process(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_page_header(slide, context, data, centered=True)
    items = _render_items(data)
    if not items:
        items = [
            {"title": text, "body": "", "asset_id": None}
            for text in _body(data)
        ]
    items = items[:5]
    if not items:
        return
    gap = 0.24
    width = (context.width - 1.36 - gap * (len(items) - 1)) / len(items)
    y = 2.30
    for index, item in enumerate(items):
        x = 0.68 + index * (width + gap)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + width / 2 - 0.26),
            Inches(1.78),
            Inches(0.52),
            Inches(0.52),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _rgb(context.theme.primary)
        circle.line.fill.background()
        _add_text(
            slide,
            x + width / 2 - 0.26,
            1.78,
            0.52,
            0.52,
            str(index + 1),
            context.theme,
            size=15,
            color="FFFFFF",
            bold=True,
            align="center",
            valign="middle",
        )
        _add_rect(
            slide,
            x,
            y,
            width,
            3.65,
            fill=context.theme.panel,
            line=context.theme.line,
            rounded=True,
        )
        _add_text(
            slide,
            x + 0.16,
            y + 0.35,
            width - 0.32,
            0.70,
            str(item.get("title") or ""),
            context.theme,
            size=18,
            color=context.theme.primary,
            bold=True,
            align="center",
            valign="middle",
        )
        body = str(item.get("body") or "").strip()
        if body:
            _add_text(
                slide,
                x + 0.18,
                y + 1.25,
                width - 0.36,
                1.75,
                body,
                context.theme,
                size=13,
                color=context.theme.text,
                align="center",
                valign="middle",
            )
        if index < len(items) - 1:
            _add_line(
                slide,
                x + width,
                y + 1.83,
                x + width + gap,
                y + 1.83,
                context.theme.secondary,
                1.8,
            )


def _render_conclusion(slide: Any, context: BuildContext, data: dict[str, Any]) -> None:
    _add_page_header(slide, context, data, centered=True)
    items = _render_items(data)
    if not items:
        items = [
            {"title": text, "body": "", "asset_id": None}
            for text in _body(data)
        ]
    items = items[:4]
    y = 2.02
    row_height = min(1.18, 4.5 / max(1, len(items)))
    for index, item in enumerate(items, start=1):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.15),
            Inches(y - 0.02),
            Inches(0.58),
            Inches(0.58),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = _rgb(context.theme.primary)
        circle.line.fill.background()
        _add_text(
            slide,
            1.18,
            y,
            0.55,
            0.55,
            str(index),
            context.theme,
            size=18,
            color="FFFFFF",
            bold=True,
            align="center",
            valign="middle",
        )
        title = str(item.get("title") or "").strip()
        body = str(item.get("body") or "").strip()
        _add_text(
            slide,
            1.98,
            y - 0.04,
            context.width - 3.15,
            0.45,
            title,
            context.theme,
            size=20,
            color=context.theme.primary,
            bold=True,
        )
        if body:
            _add_text(
                slide,
                1.98,
                y + 0.42,
                context.width - 3.15,
                0.42,
                body,
                context.theme,
                size=13,
                color=context.theme.text,
            )
        if index < len(items):
            _add_line(
                slide,
                1.98,
                y + row_height - 0.18,
                context.width - 1.18,
                y + row_height - 0.18,
                context.theme.line,
                0.8,
            )
        y += row_height


def build_presentation(
    project: dict[str, Any],
    *,
    base_dir: Path,
) -> tuple[Any | None, list[Issue], list[dict[str, Any]]]:
    issues: list[Issue] = []
    if not PPTX_AVAILABLE:
        issues.append(
            Issue(
                "error",
                "dependency.python_pptx",
                "python-pptx and Pillow are required to build editable PPTX files",
                "python module: pptx",
            )
        )
        return None, issues, []
    canvas = project.get("canvas", {})
    width = float(canvas.get("width_inches", 13.3333))
    height = float(canvas.get("height_inches", 7.5))
    presentation = Presentation()
    presentation.slide_width = Inches(width)
    presentation.slide_height = Inches(height)
    context = BuildContext(
        presentation=presentation,
        project=project,
        base_dir=base_dir.resolve(),
        width=width,
        height=height,
        theme=theme_from_project(project),
        assets={
            item["id"]: item
            for item in project.get("assets", [])
            if isinstance(item, dict) and isinstance(item.get("id"), str)
        },
        issues=issues,
        rendered_asset_ids=set(),
    )
    rendered: list[dict[str, Any]] = []
    slides = sorted(
        (item for item in project.get("slides", []) if isinstance(item, dict)),
        key=lambda item: item.get("number", 10**9),
    )
    for slide_data in slides:
        context.rendered_asset_ids = set()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(context.theme.background)
        render_type = infer_render_type(slide_data)
        renderer = {
            "cover": _render_cover,
            "section": _render_section,
            "statement": _render_bullets,
            "bullets": _render_bullets,
            "figure": _render_figure,
            "comparison": _render_comparison,
            "multi-panel": _render_multi_panel,
            "process": _render_process,
            "conclusion": _render_conclusion,
        }[render_type]
        renderer(slide, context, slide_data)
        source_asset_ids = {
            item
            for item in slide_data.get("source_asset_ids", [])
            if isinstance(item, str)
        }
        ignored_asset_ids = {
            item
            for item in _render(slide_data).get("ignored_asset_ids", [])
            if isinstance(item, str)
        }
        unused_asset_ids = sorted(
            source_asset_ids - context.rendered_asset_ids - ignored_asset_ids
        )
        for asset_id in unused_asset_ids:
            issues.append(
                Issue(
                    "warning",
                    "build.asset_unused",
                    "Declared slide evidence was not rendered; select it in "
                    "render.asset_ids or explicitly ignore it with a reason",
                    f"{slide_data.get('id', '')}:{asset_id}",
                )
            )
        number = slide_data.get("number")
        if isinstance(number, int):
            _add_page_number(slide, context, number)
        footer = _render(slide_data).get("footer")
        if isinstance(footer, str) and footer.strip():
            _add_text(
                slide,
                0.68,
                context.height - 0.42,
                context.width - 1.36,
                0.20,
                footer,
                context.theme,
                size=9,
                color=context.theme.muted,
            )
        rendered.append(
            {
                "id": slide_data.get("id"),
                "number": number,
                "render_type": render_type,
                "rendered_asset_ids": sorted(context.rendered_asset_ids),
                "unused_asset_ids": unused_asset_ids,
            }
        )
    if not slides:
        issues.append(
            Issue(
                "error",
                "build.no_slides",
                "project.json contains no slides to build",
                "slides",
            )
        )
    return presentation, issues, rendered
